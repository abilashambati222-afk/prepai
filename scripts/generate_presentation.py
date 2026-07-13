import sys
import subprocess
import os

# Auto-install python-pptx if missing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing required python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Page Setup (Widescreen 16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Custom Brand Colors
    DARK_BG = RGBColor(9, 13, 22)        # bg-brand-dark
    PRIMARY_COLOR = RGBColor(99, 102, 241) # brand-primary (indigo-500)
    ACCENT_COLOR = RGBColor(6, 182, 212)  # brand-secondary (cyan-500)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)  # slate-400
    
    # Helper: Set slide background to Dark Navy
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # Helper: Add slide title
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        return title_box

    # Slide 1: Title Slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PrepAI"
    p.font.name = 'Arial'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Placement Prep & Mock Interview Simulator"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "Decoupled MERN Platform featuring Semantic Resume Auditing & Conversational AI Simulation"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(24)
    p3.alignment = PP_ALIGN.LEFT

    # Slide Data Outline
    slides_data = [
        ("Problem Statement", [
            ("Static Preparation Portals", "Practice materials rely on generic question sets that lack personalized, role-based adaptation."),
            ("Keyword-Based Resume Checkers", "ATS filters check exact string patterns instead of semantic alignment with the job description."),
            ("High Cost of Mock Interviews", "Dynamic interview mockups with manual reviews are expensive and difficult to scale."),
            ("Lack of Detailed Feedback", "Candidates lack structured insights on specific weaknesses before actual placement interviews.")
        ]),
        ("Objectives & Scope", [
            ("Semantic Resume/ATS Auditing", "Parse candidate PDF resumes and evaluate them against custom job description skills."),
            ("Mock Placements simulator", "Conduct interactive sessions across technical, behavioral, HR, and coding prep tracks."),
            ("Instant Grading Feedback", "Score written candidate answers on correctness, completeness, and communication clarity."),
            ("Concept Analytics", "Provide polar concept mastery grids and progress timelines for historical performance logging.")
        ]),
        ("Proposed PrepAI Solution", [
            ("Profile-Based Question Generator", "AI customizes questions based on candidate target company, target role, and skill sets."),
            ("Semantic Response Grading", "Grades candidate text answers across technical accuracy, completeness, and clarity."),
            ("Simulated Speech Transcription", "Provides mic input transcription simulation to support conversational mock testing."),
            ("Printable Credentials", "Issues secure, print-friendly verification certificates for completed mock sessions.")
        ]),
        ("System Technology Stack", [
            ("Frontend Workspace", "Vite, React 19 SPA, Tailwind CSS v4, Framer Motion, Axios client, Recharts."),
            ("Backend Workspace", "Node.js, Express.js REST API gateway, JWT Security, bcryptjs salting."),
            ("Persistent Database", "MongoDB Atlas cloud document store managed via Mongoose ODM schemas."),
            ("AI Core Engine", "Google Gemini SDK Developer API Integration.")
        ]),
        ("System Block Architecture", [
            ("Client Presentation Layer", "React 19 SPA routes, private Route Guards, state context providers, and responsive dashboards."),
            ("API Router Gateway Layer", "Express routes protected by stateless JWT bearer authentication verification middleware."),
            ("AI Business Service Layer", "Quarantined Node.js helpers for Gemini prompt construction, PDF parsing, and evaluations."),
            ("Storage Layer", "MongoDB Atlas databases storing user credentials, interview sessions, and progress curves.")
        ]),
        ("Database ER Design", [
            ("User Collection Schema", "Stores login credentials, profile verification states, and academic/skill preferences."),
            ("Interview Collection Schema", "References User. Caches generated questions, user answers, and detailed grade reviews."),
            ("Career Collection Schema", "References User. Caches roadmap details, learning resources, and matching percentages."),
            ("CareerSnapshot Collection", "Tracks chronological score progressions to feed the Recharts Area charts.")
        ]),
        ("Session Security & JWT Flow", [
            ("Password Hashing", "Passwords are salted and encrypted using bcryptjs pre-save Mongoose schema hooks."),
            ("Stateless Session Authorization", "Signs a JWT upon successful registration or login and saves it in Client localStorage."),
            ("API Route Protection Middleware", "Express protect middleware intercepts requests, decodes headers, and injects user context.")
        ]),
        ("Candidate Onboarding Module", [
            ("Takewise Profile Configuration", "Collects graduation details, CGPA, target dream companies, target roles, and skill lists."),
            ("Tactile User Interface", "Tactile step-progression indicators and focus-ring forms collect onboarding preferences."),
            ("Synchronized Console Setup", "Connects user onboarding metadata to personalize the home dashboard metrics.")
        ]),
        ("Resume Ingestion & Parsing", [
            ("PDF Upload Ingestion", "Multer handles multipart PDF file uploads and buffers them to backend/uploads/."),
            ("Text Extraction Pipeline", "Extracts plain text content from PDF data buffers."),
            ("AI Structured Parsing", "Gemini parses extracted text and returns structured skills, education, and experience JSON."),
            ("Formatting Audit", "Generates an initial parsing confidence score based on layout formatting.")
        ]),
        ("Job Description Matching", [
            ("Alignment Matrix", "Compares parsed resume skills against target job description requirements."),
            ("ATS Score Gauge", "Calculates a semantic matching percentage and identifies critical skill gaps."),
            ("Study Roadmaps", "Recommends specific engineering projects and learning resources to address gaps.")
        ]),
        ("Mock Placement Simulator", [
            ("Multi-Track Sessions", "Conducts HR, Technical, Behavioral, Coding, Resume, and Project Deep-Dive interviews."),
            ("Dynamic Setup Panel", "Generates a 5-question mock session customized to target roles and companies."),
            ("Timer Constraints", "Implements time limits per question with automated submission on timeout.")
        ]),
        ("Speech Input Processing", [
            ("Voice Input Simulation", "Frontend simulates mic speech input, transcribing spoken answers in real-time."),
            ("Speech Dialogue Evaluation", "Sends transcripts to the backend and evaluates them using the same semantic models."),
            ("Tactile Wave Indicators", "Displays animated waves on the interface when speech is active.")
        ]),
        ("AI Evaluation Engine", [
            ("Technical Accuracy Criteria", "Verifies core concepts, correctness, and technical terms in candidate answers."),
            ("Completeness Index", "Checks if the answer covers all parts of the question prompt."),
            ("Clarity & Professionalism", "Evaluates communication clarity, vocabulary, and grammar."),
            ("STAR Method Verification", "Evaluates situation, task, action, and result details for behavioral questions.")
        ]),
        ("Analytics & Progression", [
            ("Progress Curves", "Recharts Area charts plot overall mock score history over time."),
            ("Subject Domain Grid", "Radar charts map candidate concept mastery across topic categories."),
            ("Category Performance Metrics", "Bar charts display average scores across Technical, HR, and Coding categories.")
        ]),
        ("Verification & Credentials", [
            ("Performance Threshold", "Mock sessions scoring >= 70% automatically issue a certificate of completion."),
            ("Tracking Verification IDs", "Generates a unique tracking verification code for each certificate."),
            ("Printable Formats", "Provides a print-friendly certificate layout with stamp seals for print/download.")
        ]),
        ("Production Cloud Hosting", [
            ("Frontend Static Bundling", "Vite React production build served on Vercel CDN edge servers."),
            ("Backend Service Hosting", "Express Node.js REST API server hosted on Render."),
            ("Cloud Database Service", "MongoDB Atlas shared cluster handles persistent storage.")
        ]),
        ("Summary & Future Scope", [
            ("Project Summary", "Developed a decoupled, AI-driven placement preparation platform with real-time semantic evaluations."),
            ("Real-Time Speech Analysis", "Future Scope: Integrate WebRTC for actual audio analysis."),
            ("Secure Sandbox Environment", "Future Scope: Add code sandbox containers to compile and run candidate code."),
            ("Multi-Tenant Administration", "Future Scope: Build dashboards for universities to track overall class readiness.")
        ])
    ]
    
    # Populate Slides
    for slide_title, bullet_points in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide)
        add_slide_header(slide, slide_title)
        
        # Left and Right layout spacing calculation
        left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.2))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        for idx, (headline, body) in enumerate(bullet_points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"▪ {headline}: "
            p.font.name = 'Arial'
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_before = Pt(14)
            
            # Append body text in white/slate
            run = p.add_run()
            run.text = body
            run.font.name = 'Arial'
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.color.rgb = TEXT_WHITE
            
    # Save PPTX
    output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation PPTX file generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
